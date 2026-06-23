from __future__ import annotations

import csv
from pathlib import Path
from typing import Dict, List

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from .models import AppError, ParsedDocument
from .storage import new_id


def _section_label(text: str) -> str:
    if any(key in text for key in ["团队", "创始", "履历", "经验"]):
        return "团队"
    if any(key in text for key in ["客户", "用户", "合同", "案例"]):
        return "客户与需求"
    if any(key in text for key in ["收入", "营收", "GMV", "回款", "利润"]):
        return "经营数据"
    if any(key in text for key in ["市场", "规模", "赛道", "行业"]):
        return "市场"
    if any(key in text for key in ["融资", "资金", "用途", "估值"]):
        return "融资"
    if any(key in text for key in ["产品", "技术", "算法", "数据", "壁垒"]):
        return "产品与技术"
    return "正文"


def make_chunks(text: str, page_texts: List[Dict[str, object]] | None = None) -> List[Dict[str, object]]:
    if page_texts is None:
        blocks = [block.strip() for block in text.replace("\r", "\n").split("\n\n") if block.strip()]
        page_texts = [{"page": index // 4 + 1, "text": block} for index, block in enumerate(blocks or [text])]

    chunks: List[Dict[str, object]] = []
    chunk_index = 0
    for page in page_texts:
        page_number = int(page["page"])
        block = str(page["text"]).strip()
        if not block:
            continue
        pieces = [block[i : i + 900] for i in range(0, len(block), 900)]
        for piece in pieces:
            chunks.append(
                {
                    "id": new_id("chunk"),
                    "chunk_index": chunk_index,
                    "page_number": page_number,
                    "section_label": _section_label(piece),
                    "text": piece.strip(),
                }
            )
            chunk_index += 1
    return chunks


def parse_text_file(path: Path) -> ParsedDocument:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return ParsedDocument(text=text, chunks=make_chunks(text), parser="text")


def parse_csv(path: Path) -> ParsedDocument:
    rows = []
    with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for index, row in enumerate(reader):
            if index >= 80:
                rows.append(["..."])
                break
            rows.append(row[:20])
    text_rows = [" | ".join(cell.strip() for cell in row) for row in rows if any(cell.strip() for cell in row)]
    text = "\n".join(text_rows)
    return ParsedDocument(text=text, chunks=make_chunks(text), parser="csv")


def parse_xlsx(path: Path) -> ParsedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise AppError("解析 XLSX 需要 openpyxl，请先安装 requirements.txt 中的依赖。") from exc

    workbook = load_workbook(path, data_only=True, read_only=True)
    pages = []
    for sheet_index, sheet in enumerate(workbook.worksheets[:8], start=1):
        lines = [f"Sheet: {sheet.title}"]
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 80:
                lines.append("...")
                break
            values = ["" if value is None else str(value).strip() for value in row[:20]]
            if any(values):
                lines.append(" | ".join(values))
        pages.append({"page": sheet_index, "text": "\n".join(lines)})
    text = "\n\n".join(str(page["text"]) for page in pages)
    return ParsedDocument(text=text, chunks=make_chunks(text, pages), parser="xlsx")


def parse_pdf(path: Path) -> ParsedDocument:
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append({"page": index, "text": page.extract_text() or ""})
    text = "\n\n".join(str(page["text"]) for page in pages)
    return ParsedDocument(text=text, chunks=make_chunks(text, pages), parser="pdf")


def parse_docx(path: Path) -> ParsedDocument:
    doc = DocxDocument(str(path))
    paragraphs = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    text = "\n\n".join(paragraphs)
    return ParsedDocument(text=text, chunks=make_chunks(text), parser="docx")


def parse_pptx(path: Path) -> ParsedDocument:
    presentation = Presentation(str(path))
    pages = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        pages.append({"page": index, "text": "\n".join(texts)})
    text = "\n\n".join(str(page["text"]) for page in pages)
    return ParsedDocument(text=text, chunks=make_chunks(text, pages), parser="pptx")


def parse_document(path: Path, pasted_text: str = "") -> ParsedDocument:
    pasted_text = pasted_text.strip()
    if pasted_text:
        return ParsedDocument(text=pasted_text, chunks=make_chunks(pasted_text), parser="pasted_text")

    suffix = path.suffix.lower()
    if suffix in [".txt", ".md", ".json"]:
        return parse_text_file(path)
    if suffix == ".csv":
        return parse_csv(path)
    if suffix == ".xlsx":
        return parse_xlsx(path)
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".docx":
        return parse_docx(path)
    if suffix in [".pptx", ".ppt"]:
        if suffix == ".ppt":
            raise AppError("第一版仅支持 PPTX；请先将 PPT 另存为 PPTX。")
        return parse_pptx(path)
    if suffix in [".png", ".jpg", ".jpeg", ".webp"]:
        raise AppError("图片 OCR 暂未接入，请粘贴 OCR 后的正文。")
    raise AppError("暂不支持该文件格式，请上传 TXT/CSV/XLSX/PDF/DOCX/PPTX 或粘贴文本。")
